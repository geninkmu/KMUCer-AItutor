
import os, io, time, json, pickle
import streamlit as st
import numpy as np
import chromadb
from chromadb.utils import embedding_functions
import pdfplumber
from pptx import Presentation
from openai import OpenAI

# ---------- Helpers ----------
def clean_text(txt: str) -> str:
    txt = txt.replace("\x00", "").replace("\ufeff", "")
    return "\n".join([line.strip() for line in txt.splitlines() if line.strip()])

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150):
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end].strip())
        if end == len(text):
            break
        start = max(0, end - overlap)
    return chunks

def extract_from_pdf_bytes(name: str, file_bytes: bytes):
    out = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                text = ""
            text = clean_text(text)
            if text:
                out.append({"source": name, "loc": f"p.{i}", "text": text})
    return out

def extract_from_pptx_bytes(name: str, file_bytes: bytes):
    out = []
    prs = Presentation(io.BytesIO(file_bytes))
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if texts:
            text = clean_text("\n".join(texts))
            if text:
                out.append({"source": name, "loc": f"slide {i}", "text": text})
    return out

def embed_texts(client, texts, model: str):
    vecs = []
    # batch to be efficient and avoid rate limits
    BATCH = 64
    for i in range(0, len(texts), BATCH):
        batch = texts[i:i+BATCH]
        resp = client.embeddings.create(model=model, input=batch)
        for d in resp.data:
            vecs.append(d.embedding)
    X = np.array(vecs, dtype="float32")
    faiss.normalize_L2(X)
    return X

def build_index(client, docs, chunk_size=800, overlap=150, model="text-embedding-3-small"):
    chunks, meta = [], []
    for rec in docs:
        for j, ch in enumerate(chunk_text(rec["text"], chunk_size, overlap)):
            chunks.append(ch)
            meta.append({"source": rec["source"], "loc": rec["loc"], "chunk_id": j})
    if not chunks:
        return None, None, None
    X = embed_texts(client, chunks, model=model)
    index = faiss.IndexFlatIP(X.shape[1])
    index.add(X)
    return index, chunks, meta

def embed_query(client, text, model="text-embedding-3-small"):
    resp = client.embeddings.create(model=model, input=[text])
    v = np.array(resp.data[0].embedding, dtype="float32")
    faiss.normalize_L2(v.reshape(1, -1))
    return v

def retrieve(client, index, store, query, top_k=4, model="text-embedding-3-small"):
    qv = embed_query(client, query, model=model)
    D, I = index.search(qv.reshape(1, -1), top_k)
    hits = []
    for rank, idx in enumerate(I[0]):
        hits.append({
            "rank": rank+1,
            "score": float(D[0][rank]),
            "text": store["chunks"][idx],
            "meta": store["meta"][idx],
        })
    return hits

def load_system_prompt():
    try:
        with open("prompts/kmucer_system_prompt.txt", "r", encoding="utf-8") as f:
            return f.read()
    except:
        return "你是 KMUCer，請用專業且親切的語氣回答學生問題。"

# ---------- Streamlit UI ----------
st.set_page_config(page_title="KMUCer — 分析化學 AI 助教", page_icon="📚", layout="wide")
st.title("📚 KMUCer — 分析化學 AI 助教（無需命令列）")
st.caption("老師上傳教材 → 一鍵建立知識庫 → 學生即可開始對話（附幽默學姊模式 & 考前衝刺）")

with st.sidebar:
    st.subheader("🔑 OpenAI 設定")
    api_key = st.text_input("OpenAI API Key", type="password", help="不會離開此頁面；只用於嵌入與回答。")
    model = st.selectbox("對話模型", ["gpt-4o", "gpt-4o-mini"], index=0)
    embed_model = st.selectbox("嵌入模型（建立索引用）", ["text-embedding-3-small", "text-embedding-3-large"], index=0)
    humor = st.toggle("幽默模式", value=True)
    exam_mode = st.toggle("考前衝刺模式（附小抄）", value=False)
    top_k = st.slider("檢索段落數", 2, 8, 4)
    temperature = st.slider("創造力 (temperature)", 0.0, 1.2, 0.7, 0.1)

if not api_key:
    st.warning("請在左側輸入 OpenAI API Key 才能使用。")
    st.stop()

client = OpenAI(api_key=api_key)
system_prompt = load_system_prompt()

tab_chat, tab_build = st.tabs(["💬 KMUCer 助教", "🧱 建立知識庫（老師用）"])

# -------- Build tab --------
with tab_build:
    st.subheader("上傳教材檔案")
    files = st.file_uploader("支援 PDF / PPTX / TXT（可多檔）", type=["pdf", "pptx", "txt", "md"], accept_multiple_files=True)
    chunk_size = st.number_input("分段大小（字數）", min_value=400, max_value=1600, value=800, step=50)
    overlap = st.number_input("分段重疊（字數）", min_value=50, max_value=400, value=150, step=10)
    build_btn = st.button("🚀 建立 / 重建索引", type="primary")

    if "store" not in st.session_state:
        st.session_state.store = None
    if "index" not in st.session_state:
        st.session_state.index = None

    if build_btn:
        if not files:
            st.error("請先上傳至少一個教材檔。")
        else:
            docs = []
            with st.spinner("抽取文字中…"):
                for f in files:
                    ext = os.path.splitext(f.name)[1].lower()
                    try:
                        if ext == ".pdf":
                            docs.extend(extract_from_pdf_bytes(f.name, f.read()))
                        elif ext == ".pptx":
                            docs.extend(extract_from_pptx_bytes(f.name, f.read()))
                        elif ext in [".txt", ".md"]:
                            txt = f.read().decode("utf-8", errors="ignore")
                            txt = clean_text(txt)
                            if txt:
                                docs.append({"source": f.name, "loc": "full", "text": txt})
                        else:
                            st.warning(f"不支援的格式：{f.name}")
                    except Exception as e:
                        st.warning(f"抽取失敗 {f.name}: {e}")

            if not docs:
                st.error("沒有抽取到任何文字，可能是掃描型 PDF（需先 OCR）或簡報是舊版 .ppt。")
            else:
                with st.spinner("建立向量索引中（需要一點時間）…"):
                    index, chunks, meta = build_index(client, docs, chunk_size, overlap, model=embed_model)
                if index is None:
                    st.error("建立索引失敗。")
                else:
                    st.session_state.index = index
                    st.session_state.store = {"chunks": chunks, "meta": meta}
                    st.success(f"索引完成！共建立 {len(chunks)} 個段落。")

    if st.session_state.store:
        st.info(f"目前索引段落數：{len(st.session_state.store['chunks'])}")

# -------- Chat tab --------
with tab_chat:
    if "index" not in st.session_state or st.session_state.index is None:
        st.warning("尚未建立知識庫。請切到『建立知識庫（老師用）』上傳教材並建立索引。")
    else:
        if "chat" not in st.session_state:
            st.session_state.chat = []

        prompt = st.chat_input("輸入你的問題（例如：為什麼 CO₂ 的對稱伸縮是 Raman active 而 IR inactive？）")
        if prompt:
            # retrieve
            blocks = retrieve(client, st.session_state.index, st.session_state.store, prompt, top_k=top_k, model=embed_model)

            # build messages
            context_text = "\n\n---\n".join(
                [f"[來源 {i+1}] ({b['meta']['source']} • {b['meta']['loc']})\n{b['text']}" for i, b in enumerate(blocks)]
            )
            style = "請保持幽默學姊風格，用簡短比喻暖場，" if humor else "保持中性專業語氣，"
            exam = "學生需要考前重點整合，請在最後附上『考前小抄』的三點要點；" if exam_mode else "最後用一行『考試小提醒』結尾；"
            sys = load_system_prompt() + "\n" + style + exam + "\n重要：只允許根據下方 Context 內容回答；若不足請如實說明。\n"
            messages = [
                {"role": "system", "content": sys},
                {"role": "user", "content": f"問題：{prompt}\n\nContext：\n{context_text}"}
            ]

            with st.spinner("KMUCer 思考中…"):
                resp = client.chat.completions.create(
                    model=model,
                    temperature=temperature,
                    messages=messages
                )
                answer = resp.choices[0].message.content

            st.session_state.chat.append({"role": "user", "content": prompt})
            st.session_state.chat.append({"role": "assistant", "content": answer, "context": blocks})

        # render
        for turn in st.session_state.get("chat", []):
            if turn["role"] == "user":
                with st.chat_message("user"):
                    st.markdown(turn["content"])
            else:
                with st.chat_message("assistant"):
                    st.markdown(turn["content"])
                    if "context" in turn:
                        with st.expander("📎 本次回答參考來源"):
                            for b in turn["context"]:
                                st.write(f"- {b['meta']['source']}（{b['meta']['loc']}）")
